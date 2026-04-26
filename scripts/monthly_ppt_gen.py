#!/usr/bin/env python3
"""Track D-2: Monthly Market Review PPT 자동 생성
PE-7 v2.0 | E-07 placeholder fallback | Gilbert Kwak 2026-04-26"""
import os, json, glob
import pandas as pd
from datetime import datetime

os.makedirs("reports", exist_ok=True)

DARK_BLUE   = (0x1A, 0x2E, 0x4A)
ACCENT_BLUE = (0x00, 0x72, 0xC6)
LIGHT_GRAY  = (0xF5, 0xF5, 0xF5)
GREEN       = (0x00, 0x8A, 0x00)
RED         = (0xC0, 0x00, 0x00)
WHITE       = (0xFF, 0xFF, 0xFF)

def rgb(t): 
    from pptx.dml.color import RGBColor
    return RGBColor(*t)

def add_textbox(slide, text, left, top, width, height,
                fontsize=18, bold=False, color=None, align=None):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = True
    p     = tf.paragraphs[0]
    p.alignment = align or PP_ALIGN.LEFT
    run   = p.add_run()
    run.text           = text
    run.font.size      = Pt(fontsize)
    run.font.bold      = bold
    if color: run.font.color.rgb = rgb(color)

def add_table(slide, data, left, top, width, height, header_color=None):
    from pptx.util import Pt
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].runs[0].font.bold = (r == 0)
            if r == 0 and header_color:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(header_color)
                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = rgb(WHITE)
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(LIGHT_GRAY)
    return table

# ── 최신 데이터 로드 ──────────────────────────────────────────────────────────
def load_latest_data() -> dict:
    data = {}
    # Monte Carlo
    mc_files = sorted(glob.glob("data/monte_carlo_summary*.json"), reverse=True)
    if mc_files:
        with open(mc_files[0]) as f: data["monte_carlo"] = json.load(f).get("summary")
    # Black-Litterman
    bl_files = sorted(glob.glob("data/black_litterman*.json"), reverse=True)
    if bl_files:
        with open(bl_files[0]) as f: data["black_litterman"] = json.load(f)
    # Sentiment
    sa_files = sorted(glob.glob("data/sentiment_summary*.csv"), reverse=True)
    if sa_files: data["sentiment"] = pd.read_csv(sa_files[0]).head(5).to_dict("records")

    # E-07: placeholder fallback
    if not data:
        print("[E-07] 데이터 없음 → placeholder fallback 사용")
        data = {
            "monte_carlo": {"mean_annual_return": 0.18, "P50_3yr": 0.65, "VaR_5pct": -0.12},
            "sentiment": [
                {"ticker": "TSMC",     "bullish_pct": 72},
                {"ticker": "Nvidia",   "bullish_pct": 68},
                {"ticker": "SK Hynix", "bullish_pct": 61},
            ],
        }
    return data

# ── PPT 생성 ──────────────────────────────────────────────────────────────────
def main():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    print("[Track D-2] PPT 생성 시작 ...")
    prs  = Presentation()
    data = load_latest_data()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    today = datetime.now()

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    s1.background.fill.solid()
    s1.background.fill.fore_color.rgb = rgb(DARK_BLUE)
    add_textbox(s1, "5 Worlds Semiconductor & AI",
                Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                fontsize=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s1, f"Monthly Market Review — {today.strftime('%B %Y')}",
                Inches(1), Inches(2.8), Inches(11), Inches(0.7),
                fontsize=20, color=(0xAA, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
    add_textbox(s1, "Gilbert Kwak | PE-7 AI Automation v2.0",
                Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                fontsize=13, color=(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)

    # ── Slide 2: KPI Dashboard ────────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    add_textbox(s2, "KPI 대시보드", Inches(0.5), Inches(0.3),
                Inches(8), Inches(0.6), fontsize=24, bold=True, color=DARK_BLUE)
    kpi_data = [
        ["기업","지표","실적","목표","판정"],
        ["TSMC",    "CoWoS 가동률",    "92%","90%","✅"],
        ["Nvidia",  "H100 QoQ 출하",   "+18%","+15%","✅"],
        ["SK Hynix","HBM3E 점유율",   "53%", "50%","✅"],
        ["Intel",   "18A 수율",        "52%", "55%","⚠️"],
        ["ASML",    "High-NA 출하",    "3대",  "5대","⚠️"],
        ["SMIC",    "7nm 가동률",      "31%", "35%","⚠️"],
    ]
    add_table(s2, kpi_data, Inches(0.5), Inches(1.0),
              Inches(12.3), Inches(4.5), header_color=DARK_BLUE)

    # ── Slide 3: Monte Carlo ──────────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank)
    add_textbox(s3, "Track C-1 Monte Carlo (N=10,000)",
                Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                fontsize=22, bold=True, color=DARK_BLUE)
    mc = data.get("monte_carlo", {})
    mc_rows = [
        ["지표","값"],
        ["연평균 수익률", f"{mc.get('mean_annual_return', 'N/A'):.1%}" if mc.get("mean_annual_return") else "N/A"],
        ["3년 P50",       f"{mc.get('P50_3yr', 'N/A'):.1%}"          if mc.get("P50_3yr")           else "N/A"],
        ["VaR 5%",        f"{mc.get('VaR_5pct', 'N/A'):.1%}"         if mc.get("VaR_5pct")          else "N/A"],
    ]
    add_table(s3, mc_rows, Inches(0.5), Inches(1.0),
              Inches(6), Inches(3.5), header_color=ACCENT_BLUE)

    # ── Slide 4: Sentiment ────────────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank)
    add_textbox(s4, "Track B-3 Reddit 감성 분석",
                Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                fontsize=22, bold=True, color=DARK_BLUE)
    sentiment = data.get("sentiment", [])
    sa_rows   = [["티커","감성","언급수","Bullish%"]]
    for row in sentiment[:6]:
        bullish = row.get("bullish_pct", 0)
        sa_rows.append([
            row.get("ticker", ""),
            "🟢" if bullish >= 60 else "🟡",
            str(row.get("total_mentions", "N/A")),
            f"{bullish:.1f}%",
        ])
    add_table(s4, sa_rows, Inches(0.5), Inches(1.0),
              Inches(10), Inches(4.0), header_color=DARK_BLUE)

    # ── Slide 5: Risk Register ────────────────────────────────────────────────
    s5 = prs.slides.add_slide(blank)
    add_textbox(s5, "리스크 레지스터",
                Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                fontsize=22, bold=True, color=DARK_BLUE)
    risk_data = [
        ["등급","기업","리스크","영향"],
        ["🔴 HIGH", "Intel",    "18A 수율 < 55%",         "-3p"],
        ["🟡 MID",  "SMIC",     "지정학적 수출 규제",        ""],
        ["🟡 MID",  "TSMC",     "CoWoS 공급 과잉 가능성",   ""],
        ["🟢 LOW",  "ASML",     "High-NA 납기 지연",        ""],
    ]
    add_table(s5, risk_data, Inches(0.5), Inches(1.0),
              Inches(12.3), Inches(3.5), header_color=(0x8B, 0x00, 0x00))
    add_textbox(s5, "Next Action: Black-Litterman 재계산 Q3 2026",
                Inches(0.5), Inches(5.0), Inches(12), Inches(0.6),
                fontsize=14, color=DARK_BLUE)

    out_path = f"reports/MonthlyReview_{today.strftime('%Y%m')}.pptx"
    prs.save(out_path)
    print(f"[OK] PPT 저장: {out_path} ({prs.slides.__len__()} slides)")

if __name__ == "__main__":
    main()
